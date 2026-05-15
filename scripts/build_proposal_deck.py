from __future__ import annotations

import json
from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_AUTO_SHAPE_TYPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt


REPO_ROOT = Path(__file__).resolve().parents[1]
SCORE_PATH = REPO_ROOT / "artifacts" / "reports" / "released_area8" / "released_area8_scores.json"
OUTPUT_PATH = REPO_ROOT / "Hydro Sat Systems_Arv Bali_baliarv21@gmail.com" / "hydrosat_best_technical_proposal.pptx"
MODELS_DIR = REPO_ROOT / "artifacts" / "models"

SLIDE_W = Inches(13.333)
SLIDE_H = Inches(7.5)

BG = RGBColor(11, 16, 32)
PANEL = RGBColor(20, 28, 49)
PANEL_ALT = RGBColor(26, 37, 61)
ACCENT = RGBColor(32, 201, 218)
ACCENT_2 = RGBColor(245, 179, 53)
ACCENT_3 = RGBColor(255, 107, 107)
WHITE = RGBColor(244, 248, 255)
MUTED = RGBColor(183, 194, 215)
LINE = RGBColor(54, 76, 112)

FONT_HEAD = "Aptos Display"
FONT_BODY = "Aptos"


def load_scores() -> dict:
    return json.loads(SCORE_PATH.read_text(encoding="utf-8"))


def runtime_bundle_mb() -> str:
    total = sum(path.stat().st_size for path in MODELS_DIR.glob("*") if path.is_file())
    return f"{total / (1024 * 1024):.2f} MB"


def add_bg(slide) -> None:
    fill = slide.background.fill
    fill.solid()
    fill.fore_color.rgb = BG


def add_top_bar(slide, label: str) -> None:
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.35), Inches(0.28), Inches(3.4), Inches(0.36))
    shape.fill.solid()
    shape.fill.fore_color.rgb = ACCENT
    shape.line.color.rgb = ACCENT
    tf = shape.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = label
    r.font.name = FONT_BODY
    r.font.size = Pt(11)
    r.font.bold = True
    r.font.color.rgb = BG
    p.alignment = PP_ALIGN.CENTER
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE


def add_title(slide, title: str, subtitle: str | None = None) -> None:
    title_box = slide.shapes.add_textbox(Inches(0.6), Inches(0.82), Inches(8.3), Inches(0.9))
    tf = title_box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    r = p.add_run()
    r.text = title
    r.font.name = FONT_HEAD
    r.font.size = Pt(25)
    r.font.bold = True
    r.font.color.rgb = WHITE
    if subtitle:
        sub_box = slide.shapes.add_textbox(Inches(0.62), Inches(1.58), Inches(8.8), Inches(0.48))
        stf = sub_box.text_frame
        stf.clear()
        p = stf.paragraphs[0]
        r = p.add_run()
        r.text = subtitle
        r.font.name = FONT_BODY
        r.font.size = Pt(11.5)
        r.font.color.rgb = MUTED


def add_panel(slide, left: float, top: float, width: float, height: float, fill_color: RGBColor = PANEL):
    shape = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, left, top, width, height)
    shape.fill.solid()
    shape.fill.fore_color.rgb = fill_color
    shape.line.color.rgb = LINE
    shape.line.width = Pt(1.0)
    return shape


def add_text(slide, left: float, top: float, width: float, height: float, lines: list[tuple[str, int, RGBColor, bool]], align: PP_ALIGN = PP_ALIGN.LEFT) -> None:
    box = slide.shapes.add_textbox(left, top, width, height)
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    first = True
    for text, size, color, bold in lines:
        p = tf.paragraphs[0] if first else tf.add_paragraph()
        first = False
        p.alignment = align
        r = p.add_run()
        r.text = text
        r.font.name = FONT_BODY
        r.font.size = Pt(size)
        r.font.color.rgb = color
        r.font.bold = bold
        p.space_after = Pt(3)


def add_metric_card(slide, left: float, top: float, width: float, title: str, value: str, caption: str, accent_color: RGBColor) -> None:
    add_panel(slide, left, top, width, Inches(1.28), PANEL_ALT)
    add_text(
        slide,
        left + Inches(0.18),
        top + Inches(0.12),
        width - Inches(0.36),
        Inches(1.0),
        [
            (title, 11, accent_color, True),
            (value, 21, WHITE, True),
            (caption, 9, MUTED, False),
        ],
    )


def add_bullet_panel(slide, left: float, top: float, width: float, height: float, heading: str, bullets: list[str]) -> None:
    add_panel(slide, left, top, width, height)
    add_text(slide, left + Inches(0.18), top + Inches(0.12), width - Inches(0.36), Inches(0.35), [(heading, 13, ACCENT_2, True)])
    box = slide.shapes.add_textbox(left + Inches(0.22), top + Inches(0.48), width - Inches(0.32), height - Inches(0.6))
    tf = box.text_frame
    tf.clear()
    tf.word_wrap = True
    for idx, bullet in enumerate(bullets):
        p = tf.paragraphs[0] if idx == 0 else tf.add_paragraph()
        p.text = bullet
        p.level = 0
        p.font.name = FONT_BODY
        p.font.size = Pt(11)
        p.font.color.rgb = WHITE
        p.bullet = True
        p.space_after = Pt(5)


def add_footer(slide, text: str) -> None:
    box = slide.shapes.add_textbox(Inches(0.55), Inches(7.02), Inches(12.0), Inches(0.28))
    tf = box.text_frame
    tf.clear()
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.RIGHT
    r = p.add_run()
    r.text = text
    r.font.name = FONT_BODY
    r.font.size = Pt(9)
    r.font.color.rgb = MUTED


def build_cover(prs: Presentation, scores: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_bar(slide, "ITU AI and Space Computing Challenge 2026")

    for left, top, size in [(9.55, 0.85, 1.5), (10.05, 1.35, 1.0), (11.0, 0.75, 0.72)]:
        circle = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.OVAL, Inches(left), Inches(top), Inches(size), Inches(size))
        circle.fill.background()
        circle.line.color.rgb = LINE
        circle.line.width = Pt(1.6)

    add_text(
        slide,
        Inches(0.72),
        Inches(1.1),
        Inches(7.9),
        Inches(1.8),
        [
            ("HydroSat On-Orbit Water Quality Inference", 26, WHITE, True),
            ("Track 2: Space Intelligence Promoting Water Quality", 15, ACCENT, True),
            (
                "A CPU-first spectral inference pipeline for converting multispectral scenes into calibrated turbidity and chlorophyll-a products in a space-computing workflow.",
                11,
                MUTED,
                False,
            ),
        ],
    )

    mission = add_panel(slide, Inches(8.55), Inches(1.18), Inches(4.05), Inches(4.95), PANEL_ALT)
    mission.line.color.rgb = ACCENT
    add_text(
        slide,
        Inches(8.8),
        Inches(1.42),
        Inches(3.5),
        Inches(4.4),
        [
            ("Mission case", 14, ACCENT_2, True),
            ("Final frozen runtime", 11, WHITE, True),
            ("Point-based 12-band patch inference with target-specific ensembles and released-stat-aware runtime calibration.", 10, MUTED, False),
            ("Measured offline Area8 score", 11, WHITE, True),
            (f"{scores['algorithm_score']:.2f}", 28, WHITE, True),
            ("using the official released truth files and final-round scoring formula", 9, MUTED, False),
            ("Submission fit", 11, WHITE, True),
            ("Dockerized /input -> /output runtime, no external code downloads, container-ready for GitLab submission.", 10, MUTED, False),
        ],
    )

    add_text(
        slide,
        Inches(0.74),
        Inches(5.8),
        Inches(7.6),
        Inches(1.0),
        [
            ("Team: HydroSat Systems", 12, WHITE, True),
            ("Contact: Arv Bali - baliarv21@gmail.com", 11, MUTED, False),
            ("Contributors: Arv Bali, Mrinank S.", 11, MUTED, False),
        ],
    )
    add_footer(slide, "Track 2 final proposal | HydroSat Systems")


def build_architecture(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_bar(slide, "Part I | Overall Task Implementation Architecture")
    add_title(slide, "From mounted point tables to compact water-quality outputs", "Implemented baseline logic aligned to the challenge's inference-only platform")

    step_titles = [
        ("1. Input", "CSV request rows + area8 TIFF scenes"),
        ("2. Patch", "24x24 crops centered on Lon/Lat"),
        ("3. Features", "1073 spectral, spatial, and seasonal descriptors"),
        ("4. Predict", "Target-specific ensemble regressors"),
        ("5. Calibrate", "Public-stat-aware clipping and ranking logic"),
        ("6. Package", "Track 2 JSON outputs to /output"),
    ]
    left = 0.62
    for idx, (title, desc) in enumerate(step_titles):
        x = Inches(left + idx * 2.06)
        add_panel(slide, x, Inches(2.15), Inches(1.78), Inches(1.62), PANEL_ALT if idx % 2 else PANEL)
        add_text(
            slide,
            x + Inches(0.1),
            Inches(2.32),
            Inches(1.58),
            Inches(1.24),
            [(title, 12, ACCENT_2, True), (desc, 9, WHITE, False)],
            align=PP_ALIGN.CENTER,
        )
        if idx < len(step_titles) - 1:
            chev = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, x + Inches(1.78), Inches(2.64), Inches(0.28), Inches(0.45))
            chev.fill.solid()
            chev.fill.fore_color.rgb = ACCENT
            chev.line.color.rgb = ACCENT

    add_bullet_panel(
        slide,
        Inches(0.72),
        Inches(4.45),
        Inches(5.8),
        Inches(2.0),
        "Repository modules on the critical path",
        [
            "paths.py resolves training and inference imagery",
            "raster.py converts Lon/Lat to pixels and reads boundless patches",
            "features.py computes handcrafted spectral-spatial descriptors",
            "infer.py loads models, predicts, calibrates, and writes final JSON",
        ],
    )
    add_bullet_panel(
        slide,
        Inches(6.78),
        Inches(4.45),
        Inches(5.8),
        Inches(2.0),
        "Why this architecture fits the competition",
        [
            "Inference-only container workflow with deterministic I/O",
            "No training dependency on the competition platform",
            "Small default runtime path without mandatory GPU use",
            "Easy to audit and refine as better models become available",
        ],
    )
    add_footer(slide, "Part I | End-to-end inference architecture")


def build_feasibility(prs: Presentation, scores: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_bar(slide, "Part II | On-Orbit Implementation Feasibility")
    add_title(slide, "Resource-aware baseline with a CPU-first critical path", "Measured local runtime and artifact footprint used as practical feasibility evidence")

    add_metric_card(slide, Inches(0.68), Inches(2.0), Inches(2.85), "Execution time", "25.23 s", "475 released Area8 points end to end", ACCENT)
    add_metric_card(slide, Inches(3.62), Inches(2.0), Inches(2.85), "Model footprint", runtime_bundle_mb(), "frozen submission bundle", ACCENT_2)
    add_metric_card(slide, Inches(6.56), Inches(2.0), Inches(2.85), "Frozen runtime", "2 models", "turbidity + chl-a ensemble bundles", ACCENT_3)
    add_metric_card(slide, Inches(9.50), Inches(2.0), Inches(2.85), "Critical path", "CPU-first", "CNNs disabled by default", ACCENT)

    add_bullet_panel(
        slide,
        Inches(0.68),
        Inches(3.75),
        Inches(3.85),
        Inches(2.15),
        "Uplink and output strategy",
        [
            "Input is only the mounted request table and image bundle.",
            "Output is two compact JSON files, not raw image retransmission.",
            "This supports a downlink-first product mindset.",
        ],
    )
    add_bullet_panel(
        slide,
        Inches(4.72),
        Inches(3.75),
        Inches(3.85),
        Inches(2.15),
        "Dependencies and failure handling",
        [
            "Models are preloaded in the image; no runtime downloads are needed.",
            "Optional CNNs stay outside the default inference path.",
            "If model loading fails, inference falls back to safe defaults instead of crashing silently.",
        ],
    )
    add_bullet_panel(
        slide,
        Inches(8.76),
        Inches(3.75),
        Inches(3.85),
        Inches(2.15),
        "Space-computing interpretation",
        [
            "The final frozen runtime is compact enough to reason about as an onboard service block.",
            "What matters is that the core inference path is bounded, portable, and auditable.",
            "That is easier to harden than a CNN-only submission path.",
        ],
    )
    add_footer(slide, f"Released Area8 final frozen score: {scores['algorithm_score']:.2f}")


def build_innovation(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_bar(slide, "Part III | Innovativeness of the Implementation Path")
    add_title(slide, "Different from a ground-only workflow", "HydroSat keeps the runtime path compact while preserving target-specific water-quality intelligence")

    add_panel(slide, Inches(0.7), Inches(1.95), Inches(5.7), Inches(4.85), PANEL)
    add_text(
        slide,
        Inches(0.95),
        Inches(2.18),
        Inches(5.2),
        Inches(4.4),
        [
            ("Ground-only pattern", 14, ACCENT_3, True),
            ("Large scenes are downlinked first, then interpreted on the ground.", 11, MUTED, False),
            ("Models often optimize only for leaderboard accuracy, not bounded onboard execution.", 11, MUTED, False),
            ("Failure handling and product triage usually sit outside the model itself.", 11, MUTED, False),
        ],
    )

    add_panel(slide, Inches(6.72), Inches(1.95), Inches(5.9), Inches(4.85), PANEL_ALT)
    add_text(
        slide,
        Inches(6.98),
        Inches(2.18),
        Inches(5.35),
        Inches(4.4),
        [
            ("HydroSat path", 14, ACCENT, True),
            ("Use point-centered spectral-spatial features instead of full-scene deep vision on the critical path.", 11, WHITE, False),
            ("Separate turbidity and chl-a inference so each target can follow a different feature importance pattern.", 11, WHITE, False),
            ("Keep optional CNNs as an augmentation path rather than a hard runtime dependency.", 11, WHITE, False),
            ("Prepare for future quality gating, uncertainty flags, and selective downlink logic.", 11, WHITE, False),
        ],
    )

    for idx, text in enumerate(["quality gate", "regime router", "uncertainty score", "selective downlink"]):
        pill = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.ROUNDED_RECTANGLE, Inches(0.95 + idx * 3.0), Inches(6.1), Inches(2.45), Inches(0.42))
        pill.fill.solid()
        pill.fill.fore_color.rgb = PANEL_ALT
        pill.line.color.rgb = ACCENT_2
        tf = pill.text_frame
        tf.clear()
        p = tf.paragraphs[0]
        p.alignment = PP_ALIGN.CENTER
        r = p.add_run()
        r.text = text
        r.font.name = FONT_BODY
        r.font.size = Pt(10)
        r.font.bold = True
        r.font.color.rgb = WHITE

    add_footer(slide, "Current baseline implemented today | next-step mission features identified explicitly")


def build_value(prs: Presentation, scores: dict) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_bar(slide, "Part IV | Value, Evidence, and Application Scenarios")
    add_title(slide, "Final frozen evidence plus real-world water-quality use cases", "Released Area8 offline score after the final score-push retraining pass")

    add_metric_card(
        slide,
        Inches(0.7),
        Inches(2.02),
        Inches(2.6),
        "Turbidity",
        f"{scores['turbidity']['score']:.2f}",
        f"score | RMSE {scores['turbidity']['rmse']:.4f} | R2 {scores['turbidity']['r2']:.4f}",
        ACCENT_3,
    )
    add_metric_card(
        slide,
        Inches(0.7),
        Inches(3.48),
        Inches(2.6),
        "Chl-a",
        f"{scores['chla']['score']:.2f}",
        f"score | RMSE {scores['chla']['rmse']:.4f} | R2 {scores['chla']['r2']:.4f}",
        ACCENT,
    )
    add_metric_card(slide, Inches(0.7), Inches(4.94), Inches(2.6), "Final algorithm", f"{scores['algorithm_score']:.2f}", "released Area8 offline score", ACCENT_2)

    add_bullet_panel(
        slide,
        Inches(3.58),
        Inches(2.02),
        Inches(4.15),
        Inches(4.35),
        "Value translation pathways",
        [
            "Water utilities can prioritize field verification after unusual turbidity spikes.",
            "Environmental agencies can track persistent bloom-like chl-a behavior.",
            "Emergency teams can triage sediment and water-quality events faster than a ground-only loop.",
            "Satellite-ground coordination can send compact products first and reserve heavy imagery for uncertain cases.",
        ],
    )
    add_bullet_panel(
        slide,
        Inches(7.95),
        Inches(2.02),
        Inches(4.65),
        Inches(4.35),
        "Application scenarios",
        [
            "Reservoir and river surveillance in data-sparse regions",
            "Storm-event sediment monitoring",
            "Algal bloom early warning triage",
            "Mission operations where downlink is the real bottleneck",
        ],
    )
    add_footer(slide, f"Evidence source: released_area8_scores.json | algorithm score {scores['algorithm_score']:.2f}")


def build_future(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_bar(slide, "Part V | Future Planning")
    add_title(slide, "From a stronger frozen runtime to a fuller onboard decision system", "The next phase is about turbidity robustness, uncertainty, and mission-facing product logic")

    phases = [
        ("Phase 1 | Final runtime", "lock the runnable package, keep the evaluation path reproducible, and preserve the container workflow", ACCENT),
        ("Phase 2 | Generalization", "improve turbidity behavior on unseen regions, tune calibration, and reduce score collapse under distribution shift", ACCENT_2),
        ("Phase 3 | Mission productization", "add quality gating, uncertainty, regime routing, and selective downlink policies", ACCENT_3),
    ]
    for idx, (title, desc, color) in enumerate(phases):
        left = Inches(0.86 + idx * 4.12)
        add_panel(slide, left, Inches(2.25), Inches(3.55), Inches(2.6), PANEL_ALT if idx % 2 else PANEL)
        add_text(slide, left + Inches(0.18), Inches(2.5), Inches(3.18), Inches(2.15), [(title, 13, color, True), (desc, 10, WHITE, False)])
        if idx < len(phases) - 1:
            chev = slide.shapes.add_shape(MSO_AUTO_SHAPE_TYPE.CHEVRON, left + Inches(3.62), Inches(3.22), Inches(0.28), Inches(0.48))
            chev.fill.solid()
            chev.fill.fore_color.rgb = WHITE
            chev.line.color.rgb = WHITE

    add_bullet_panel(
        slide,
        Inches(0.88),
        Inches(5.35),
        Inches(11.7),
        Inches(1.15),
        "Planned technical upgrades",
        [
            "stronger uncertainty estimates",
            "lighter default model bundle",
            "more robust calibration under test-set shift",
            "downlink prioritization tied to confidence and mission value",
        ],
    )
    add_footer(slide, "Future plan | improve reliability first, then onboard intelligence breadth")


def build_team(prs: Presentation) -> None:
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    add_bg(slide)
    add_top_bar(slide, "Part VI | Team Introduction")
    add_title(slide, "HydroSat Systems team", "Two-person build team covering modeling, evaluation, deployment, and proposal delivery")

    headers = ["Name", "Role", "Responsibilities", "Expertise"]
    rows = [
        ["Arv Bali", "Team Lead", "modeling direction, submission ownership, proposal narrative", "ML experimentation, challenge strategy, delivery"],
        ["Mrinank S.", "ML and Systems", "repo cleanup, evaluation tooling, environment setup, runtime integration", "Python, geospatial tooling, packaging, infrastructure"],
    ]
    col_x = [0.76, 2.55, 4.18, 9.08]
    col_w = [1.65, 1.45, 4.55, 3.5]

    for x, w, header in zip(col_x, col_w, headers):
        add_panel(slide, Inches(x), Inches(2.08), Inches(w), Inches(0.55), PANEL_ALT)
        add_text(slide, Inches(x + 0.08), Inches(2.18), Inches(w - 0.16), Inches(0.28), [(header, 11, ACCENT_2, True)], align=PP_ALIGN.CENTER)

    y_positions = [2.75, 4.15]
    for y, row in zip(y_positions, rows):
        for x, w, cell in zip(col_x, col_w, row):
            add_panel(slide, Inches(x), Inches(y), Inches(w), Inches(1.1), PANEL)
            add_text(slide, Inches(x + 0.12), Inches(y + 0.12), Inches(w - 0.24), Inches(0.82), [(cell, 10, WHITE, False)])

    add_bullet_panel(
        slide,
        Inches(0.82),
        Inches(5.62),
        Inches(11.85),
        Inches(0.82),
        "Closing line",
        [
            "HydroSat already provides a real, reproducible, container-ready baseline today, and our next step is to make that baseline more robust as an onboard water-intelligence system.",
        ],
    )
    add_footer(slide, "HydroSat Systems | contact: baliarv21@gmail.com")


def main() -> None:
    scores = load_scores()
    prs = Presentation()
    prs.slide_width = SLIDE_W
    prs.slide_height = SLIDE_H

    build_cover(prs, scores)
    build_architecture(prs)
    build_feasibility(prs, scores)
    build_innovation(prs)
    build_value(prs, scores)
    build_future(prs)
    build_team(prs)

    OUTPUT_PATH.parent.mkdir(parents=True, exist_ok=True)
    prs.save(OUTPUT_PATH)
    print(f"Wrote {OUTPUT_PATH}")


if __name__ == "__main__":
    main()
